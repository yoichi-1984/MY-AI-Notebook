import io
import os
import re
import glob
import uuid
import tempfile
import hashlib
from typing import List, Dict, Any, Tuple, Optional

# 各種パーサライブラリのインポート（利用できない場合のフォールバックを考慮）
try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    import pypdfium2 as pdfium
    HAS_PYPDFIUM2 = True
except ImportError:
    HAS_PYPDFIUM2 = False

try:
    import pptx
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

# win32com (Windows専用) のダイナミックインポート用フラグ
HAS_WIN32 = False
try:
    if os.name == 'nt':
        import win32com.client
        import pythoncom
        HAS_WIN32 = True
except ImportError:
    pass

class DocumentParser:
    """PDF, PPTX, DOCX, XLSX, TXTドキュメントからテキストと画像をページ単位で抽出するパーサクラス"""

    @classmethod
    def parse(cls, file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
        """ファイル名から拡張子を判別し、適切なパース処理を実行する"""
        ext = os.path.splitext(filename)[1].lower()
        
        if ext == '.pdf':
            return cls.parse_pdf(file_bytes, filename)
        elif ext in ['.pptx', '.ppt']:
            return cls.parse_pptx(file_bytes, filename)
        elif ext == '.docx':
            return cls.parse_docx(file_bytes)
        elif ext in ['.xlsx', '.xlsm', '.xls']:
            return cls.parse_xlsx(file_bytes, filename)
        elif ext in ['.txt', '.md']:
            return cls.parse_text(file_bytes, filename)
        else:
            raise ValueError(f"未サポートのファイル形式です: {ext}")

    @classmethod
    def parse_pdf(cls, file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
        """PDFファイルを解析し、ページ単位のテキストと画像を抽出する"""
        pages = []
        if not HAS_PYPDF:
            raise ImportError("pypdf がインストールされていません。")
            
        pdf_file = io.BytesIO(file_bytes)
        reader = pypdf.PdfReader(pdf_file)
        num_pages = len(reader.pages)

        # pypdfium2によるレンダリング準備
        doc_ium = None
        if HAS_PYPDFIUM2:
            try:
                doc_ium = pdfium.PdfDocument(io.BytesIO(file_bytes))
            except Exception as e:
                print(f"[Warning] pypdfium2でPDFを開けませんでした: {e}")

        for i in range(num_pages):
            page_name = f"ページ {i + 1}"
            
            # テキスト抽出
            text = ""
            try:
                text = reader.pages[i].extract_text() or ""
            except Exception as e:
                print(f"[Warning] pypdfでのテキスト抽出に失敗しました (ページ {i+1}): {e}")

            # 画像レンダリング (プレビュー用 ＆ スキャンPDFの場合はOCR用)
            img_bytes = None
            if doc_ium:
                try:
                    page_ium = doc_ium[i]
                    # scale=2 で高画質レンダリング (約 150 DPI)
                    bitmap = page_ium.render(scale=2)
                    pil_img = bitmap.to_pil()
                    img_byte_arr = io.BytesIO()
                    pil_img.save(img_byte_arr, format='PNG')
                    img_bytes = img_byte_arr.getvalue()
                except Exception as e:
                    print(f"[Warning] pypdfium2でのレンダリングに失敗しました (ページ {i+1}): {e}")

            pages.append({
                "page_name": page_name,
                "text": text,
                "image_bytes": img_bytes,
                "image_ext": ".png" if img_bytes else None
            })

        if doc_ium:
            doc_ium.close()

        return pages

    @classmethod
    def parse_pptx(cls, file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
        """PPTXファイルを解析し、スライド単位のテキストと画像を抽出する"""
        if not HAS_PPTX:
            raise ImportError("python-pptx がインストールされていません。")

        pages = []
        
        # 1. まず python-pptx で各スライドのテキストと内部画像を抽出
        prs = Presentation(io.BytesIO(file_bytes))
        slide_texts = []
        slide_embedded_images = []

        for idx, slide in enumerate(prs.slides):
            texts = []
            # プレースホルダーやテキストボックスなどからテキストを抽出
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text.strip())
                # テーブルのテキストを抽出
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
            
            slide_texts.append("\n".join(texts))

            # 内部に埋め込まれた画像の抽出（フォールバック用）
            embedded_img = None
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # 最初の画像をスライドの代表画像として抽出
                        embedded_img = shape.image.blob
                        break
                    except Exception:
                        pass
            slide_embedded_images.append(embedded_img)

        # 2. Windows + PowerPoint環境であれば、win32comでスライドを高画質一括エクスポート
        win32_images = []
        if HAS_WIN32:
            try:
                win32_images = cls._export_ppt_slides_win32(file_bytes, filename)
            except Exception as e:
                print(f"[Warning] win32comによるPowerPoint画像化に失敗しました: {e}")

        # 3. 情報をマージして返す
        for idx in range(len(prs.slides)):
            page_name = f"スライド {idx + 1}"
            text = slide_texts[idx]
            
            # 高品質な画像があればそれを使用、なければフォールバック画像を使用
            img_bytes = None
            if idx < len(win32_images):
                img_bytes = win32_images[idx]
            elif idx < len(slide_embedded_images) and slide_embedded_images[idx]:
                img_bytes = slide_embedded_images[idx]

            pages.append({
                "page_name": page_name,
                "text": text,
                "image_bytes": img_bytes,
                "image_ext": ".png" if img_bytes else None
            })

        return pages

    @classmethod
    def _export_ppt_slides_win32(cls, file_bytes: bytes, filename: str) -> List[bytes]:
        """win32comを利用してPowerPointスライドをPNG画像としてエクスポートする"""
        if not HAS_WIN32:
            return []

        image_list = []
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_ppt_path = os.path.join(temp_dir, filename)
            with open(temp_ppt_path, "wb") as f:
                f.write(file_bytes)
            
            output_dir = os.path.join(temp_dir, "slides")
            os.makedirs(output_dir, exist_ok=True)

            ppt_app = None
            presentation = None
            
            try:
                # COMの初期化
                pythoncom.CoInitialize()
                # PowerPoint起動
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                # ファイルを開く (ReadOnly=True, WithWindow=False)
                presentation = ppt_app.Presentations.Open(
                    os.path.abspath(temp_ppt_path),
                    ReadOnly=True,
                    WithWindow=False
                )
                
                # PNG形式でエクスポート (18 = ppSaveAsPNG)
                export_path = os.path.abspath(os.path.join(output_dir, "slide.png"))
                presentation.SaveAs(export_path, 18)
            except Exception as e:
                print(f"[Error] PPTX SaveAs COM Error: {e}")
                return []
            finally:
                if presentation:
                    try:
                        presentation.Close()
                    except Exception:
                        pass
                if ppt_app:
                    try:
                        # PowerPointを閉じる
                        ppt_app.Quit()
                    except Exception:
                        pass
                ppt_app = None
                # COMの解放
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

            # エクスポートされた画像ファイルの探索とロード
            # PowerPointは "slide/slide1.png" のようにフォルダを作成する場合と、
            # "slide1.png" のように直接出力する場合があります。
            slide_files = glob.glob(os.path.join(output_dir, "*.PNG"))
            if not slide_files:
                slide_files = glob.glob(os.path.join(output_dir, "*.png"))
            if not slide_files and os.path.isdir(os.path.join(output_dir, "slide")):
                slide_files = glob.glob(os.path.join(output_dir, "slide", "*.PNG"))
                if not slide_files:
                    slide_files = glob.glob(os.path.join(output_dir, "slide", "*.png"))

            # ファイル名に含まれる数字でソートする (自然順ソート)
            def extract_number(filepath):
                nums = re.findall(r'\d+', os.path.basename(filepath))
                return int(nums[0]) if nums else 0

            slide_files.sort(key=extract_number)

            for slide_file in slide_files:
                with open(slide_file, "rb") as img_f:
                    image_list.append(img_f.read())
            
        return image_list

    @classmethod
    def parse_docx(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """DOCXファイルを解析し、見出し単位で分割したテキストを抽出する"""
        if not HAS_DOCX:
            raise ImportError("python-docx がインストールされていません。")

        doc = docx.Document(io.BytesIO(file_bytes))
        sections = []
        
        current_title = "文書の始まり"
        current_text = []

        # Word文書内の段落を順番にスキャン
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # 見出しスタイルを判定 (Heading 1 〜 Heading 3 など)
            is_heading = False
            if para.style.name.startswith("Heading") or para.style.name.startswith("見出し"):
                is_heading = True
            elif re.match(r'^第[一二三四五六七八九十\d]+章|^\d+\.\s|^\d+\.\d+\s', text):
                # テキストパターンから見出しを推測
                is_heading = True

            if is_heading:
                # 既存セクションを保存
                if current_text:
                    sections.append({
                        "page_name": current_title,
                        "text": "\n".join(current_text),
                        "image_bytes": None,
                        "image_ext": None
                    })
                current_title = text
                current_text = []
            else:
                current_text.append(text)

        # 残りの段落を保存
        if current_text or current_title != "文書の始まり":
            sections.append({
                "page_name": current_title,
                "text": "\n".join(current_text),
                "image_bytes": None,
                "image_ext": None
            })

        # 見出し分割が少なすぎる、または見出しが全くなかった場合は2000文字でチャンク分割
        if not sections or (len(sections) == 1 and len(sections[0]["text"]) > 3000):
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            return cls._chunk_text(full_text, chunk_size=2000, prefix="セクション")

        return sections

    @classmethod
    def parse_xlsx(cls, file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
        """Excelファイルを解析し、シート単位でMarkdownテーブル化したテキストを抽出する"""
        if not HAS_PANDAS:
            raise ImportError("pandas がインストールされていません。")

        pages = []
        excel_file = io.BytesIO(file_bytes)
        
        # エンジン選定
        engine = "openpyxl"
        
        try:
            # 全シートを一度に読み込む
            sheets = pd.read_excel(excel_file, sheet_name=None, engine=engine)
        except Exception as e:
            raise ValueError(f"Excelファイルの読み込みに失敗しました: {e}")

        for sheet_name, df in sheets.items():
            page_name = f"シート: {sheet_name}"
            
            if df.empty:
                text = "(空のシートです)\n"
            else:
                df_clean = df.fillna("")
                
                # tabulate が使える場合は Markdown テーブルに変換
                try:
                    text = df_clean.to_markdown(index=False)
                except ImportError:
                    # tabulateがない場合はCSVにフォールバック
                    text = df_clean.to_csv(index=False)
            
            pages.append({
                "page_name": page_name,
                "text": text,
                "image_bytes": None,
                "image_ext": None
            })

        return pages

    @classmethod
    def parse_text(cls, file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
        """テキストファイルをデコードし、見出しまたは3000文字単位で分割したものを抽出する"""
        # 文字コードのデコード試行
        text_content = ""
        for encoding in ["utf-8", "cp932", "euc-jp", "utf-16"]:
            try:
                text_content = file_bytes.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text_content = file_bytes.decode("utf-8", errors="replace")

        # Markdown見出し (# など) が含まれるかチェック
        # 見出しがあれば、見出し単位で分割
        headings = list(re.finditer(r'^(#+\s+.+)$', text_content, re.MULTILINE))
        
        if len(headings) >= 2:
            pages = []
            prev_pos = 0
            prev_title = "はじめに"
            
            for match in headings:
                start, end = match.span()
                section_text = text_content[prev_pos:start].strip()
                if section_text:
                    pages.append({
                        "page_name": prev_title,
                        "text": section_text,
                        "image_bytes": None,
                        "image_ext": None
                    })
                prev_title = match.group(1).strip().replace("#", "").strip()
                prev_pos = end

            # 最後のセクション
            section_text = text_content[prev_pos:].strip()
            if section_text:
                pages.append({
                    "page_name": prev_title,
                    "text": section_text,
                    "image_bytes": None,
                    "image_ext": None
                })
            return pages
        else:
            # 見出しがあまりない場合は3000文字ごとに分割
            return cls._chunk_text(text_content, chunk_size=3000, prefix="ページ")

    @classmethod
    def _chunk_text(cls, text: str, chunk_size: int, prefix: str) -> List[Dict[str, Any]]:
        """テキストを文字数ベースでチャンク分割する"""
        pages = []
        text_len = len(text)
        
        if text_len == 0:
            return [{
                "page_name": f"{prefix} 1",
                "text": "",
                "image_bytes": None,
                "image_ext": None
            }]

        idx = 1
        for start_pos in range(0, text_len, chunk_size):
            chunk = text[start_pos:start_pos + chunk_size].strip()
            if chunk:
                pages.append({
                    "page_name": f"{prefix} {idx}",
                    "text": chunk,
                    "image_bytes": None,
                    "image_ext": None
                })
                idx += 1
        return pages
